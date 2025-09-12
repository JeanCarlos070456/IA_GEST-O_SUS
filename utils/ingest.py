import os, json
import pandas as pd
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
from typing import Dict, Tuple, List
from unidecode import unidecode

def read_text_from_file(path: str) -> Tuple[str, Dict]:
    ext = os.path.splitext(path)[1].lower()
    meta = {"source_path": path}

    if ext == '.pdf':
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                texts.append(page.extract_text() or '')
        content = "\n".join(texts)

    elif ext == '.docx':
        doc = DocxDocument(path)
        content = "\n".join(p.text for p in doc.paragraphs)

    elif ext == '.pptx':
        prs = Presentation(path)
        slides = []
        for s in prs.slides:
            sbuf = []
            for shp in s.shapes:
                if hasattr(shp, 'text'):
                    sbuf.append(shp.text)
            slides.append("\n".join(sbuf))
        content = "\n\n".join(slides)

    elif ext in ['.csv', '.xlsx', '.xls']:
        try:
            if ext == '.csv':
                df = pd.read_csv(path, dtype=str)
            else:
                df = pd.read_excel(path, dtype=str)
        except Exception:
            df = pd.read_excel(path, engine='openpyxl', dtype=str)
        content = df.fillna('').to_csv(index=False)

    elif ext in ['.html', '.htm']:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'lxml')
            content = soup.get_text(separator='\n')

    else:  # .txt
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

    # Limpeza básica
    content = content.replace('\u200b', ' ').replace('\xa0', ' ')
    content = unidecode(content)

    return content, meta

def chunk_text(text: str, size: int, overlap: int) -> List[str]:
    text = " ".join(text.split())
    chunks, i = [], 0
    while i < len(text):
        chunks.append(text[i:i+size])
        i += max(size - overlap, 1)
    return [c for c in chunks if c.strip()]

def write_jsonl(path: str, rows: List[Dict]):
    with open(path, 'w', encoding='utf-8') as f:
        for r in rows:
            f.write(json.dumps(r, ensure_ascii=False) + "\n")
